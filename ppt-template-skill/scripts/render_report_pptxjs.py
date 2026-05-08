#!/usr/bin/env python3
from __future__ import annotations

import json
import random
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

from ppt_engine.content import (
    ensure_toc_slide,
    render_toc_slide,
    render_transition_slide,
    replace_cover_text,
)
from ppt_engine.context import RenderContext
from ppt_engine.layout_registry import LayoutRegistry
from ppt_engine.template import (
    collect_slide_fonts,
    delete_slide,
    duplicate_slide,
    load_template_tokens,
    move_slide,
    prepare_content_template,
    resolve_template_roles,
)


def matching_template_spec(template_path: Path) -> dict:
    for spec_path in template_path.parent.glob("*_spec_corrected.json"):
        try:
            spec = json.loads(spec_path.read_text(encoding="utf-8-sig"))
        except Exception:
            continue
        if spec.get("template_file") == template_path.name:
            return spec
    return {}


def default_body_zone(content_top: float | None) -> dict:
    y = max(1.3, (content_top or 1.1) + 0.28)
    return {"x": 0.62, "y": y, "w": 12.05, "h": max(1.0, 7.5 - y - 0.55)}


def slide_xml_index(slide) -> int:
    partname = str(slide.part.partname)
    match = re.search(r"slide(\d+)\.xml$", partname)
    if not match:
        raise ValueError(f"Cannot resolve slide XML index from partname: {partname}")
    return int(match.group(1))


def unzip_pptx(pptx_path: Path, output_dir: Path):
    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(pptx_path) as zf:
        zf.extractall(output_dir)


def build_skeleton_deck(content_path: Path, skeleton_pptx: Path, template_path: Path) -> dict:
    content = json.loads(content_path.read_text(encoding="utf-8-sig"))
    prs = Presentation(str(template_path))
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    registry = LayoutRegistry(skill_dir)
    ctx = RenderContext(skill_dir=skill_dir, template_path=template_path, tokens=load_template_tokens(template_path))

    cover_idx, toc_template_idx, transition_indices, content_indices, end_idx = resolve_template_roles(prs, template_path)
    inferred_fonts = collect_slide_fonts([prs.slides[idx] for idx in content_indices if idx < len(prs.slides)])
    if inferred_fonts and not ctx.tokens.get("font_body"):
        ctx.tokens["font_body"] = inferred_fonts.most_common(1)[0][0]
    if inferred_fonts and not ctx.tokens.get("font_title"):
        ctx.tokens["font_title"] = ctx.tokens.get("font_body")

    desired_slides = ensure_toc_slide(content.get("slides") or [], toc_template_idx is not None)
    body_slides = [s for s in desired_slides if s.get("type") in ("toc", "transition", "chapter", "content")]
    toc_source = prs.slides[toc_template_idx] if toc_template_idx is not None else prs.slides[content_indices[0]]
    transition_sources = [prs.slides[idx] for idx in transition_indices] or [prs.slides[content_indices[0]]]
    content_sources = [prs.slides[idx] for idx in content_indices]
    end_source = prs.slides[end_idx]
    cover_slide_id = prs.slides[cover_idx].slide_id
    end_slide_id = end_source.slide_id

    style_selection = content.get("style_selection", "rotate")
    rng = random.Random(content.get("style_seed", 20260430))
    generated = []
    content_style_cursor = 0
    transition_style_cursor = 0
    for slide_data in body_slides:
        slide_type = slide_data.get("type")
        if slide_type == "toc":
            generated.append((duplicate_slide(prs, toc_source), slide_data))
        elif slide_type in ("transition", "chapter"):
            source = transition_sources[transition_style_cursor % len(transition_sources)]
            transition_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))
        else:
            source = rng.choice(content_sources) if style_selection == "random" else content_sources[content_style_cursor % len(content_sources)]
            content_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))

    keep_ids = {cover_slide_id, end_slide_id}
    keep_ids.update(slide.slide_id for slide, _ in generated)
    for idx in range(len(prs.slides) - 1, -1, -1):
        if prs.slides[idx].slide_id not in keep_ids:
            delete_slide(prs, idx)

    target_order = [cover_slide_id] + [slide.slide_id for slide, _ in generated] + [end_slide_id]
    for target_idx, slide_id in enumerate(target_order):
        current_idx = next(i for i, slide in enumerate(prs.slides) if slide.slide_id == slide_id)
        if current_idx != target_idx:
            move_slide(prs, current_idx, target_idx)

    slides = list(prs.slides)
    title_bottom = {}
    for slide, slide_data in generated:
        if slide_data.get("type") in ("toc", "transition", "chapter"):
            continue
        title_bottom[slide.slide_id] = prepare_content_template(slide, slide_data, registry)

    cover_data = next((s for s in desired_slides if s.get("type") == "cover"), {})
    replace_cover_text(slides[0], cover_data, content, ctx)

    source = content.get("source", "")
    content_manifest_slides = []
    spec = matching_template_spec(template_path)
    spec_zone = (((spec.get("slide_types") or {}).get("content") or {}).get("body_zone") or {})
    for i, slide_data in enumerate(body_slides, start=1):
        slide = slides[i]
        slide_num = i + 1
        if slide_data.get("type") == "toc":
            render_toc_slide(slide, slide_data, slide_num, source, ctx)
        elif slide_data.get("type") in ("transition", "chapter"):
            render_transition_slide(slide, slide_data, slide_num, source, ctx)
        else:
            zone = dict(spec_zone) if spec_zone else default_body_zone(title_bottom.get(slide.slide_id))
            if title_bottom.get(slide.slide_id) is not None:
                zone["y"] = max(float(zone.get("y", 1.3)), float(title_bottom[slide.slide_id]) + 0.28)
                zone["h"] = max(1.0, 7.5 - zone["y"] - 0.55)
            content_manifest_slides.append({
                "idx": slide_xml_index(slide),
                "page_num": slide_num,
                "type": "content",
                "is_content": True,
                "body_zone": zone,
                "page_title": slide_data.get("page_title") or slide_data.get("title") or "",
                "body": slide_data.get("body") or {},
            })

    skeleton_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(skeleton_pptx))
    return {
        "content": content,
        "tokens": ctx.tokens,
        "content_slides": content_manifest_slides,
    }


def build(content_path: Path, output_path: Path, template_path: Path):
    script_dir = Path(__file__).resolve().parent
    work_dir = output_path.with_suffix("")
    work_dir = work_dir.parent / f"{work_dir.name}_pptxjs_work"
    skeleton_pptx = work_dir.parent / f"{work_dir.name}_skeleton.pptx"

    data = build_skeleton_deck(content_path, skeleton_pptx, template_path)
    unzip_pptx(skeleton_pptx, work_dir)

    manifest = {
        "output_dir": str(work_dir),
        "design_tokens": data["tokens"],
        "content_slides": data["content_slides"],
    }
    manifest_path = work_dir / "content_manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    subprocess.run(["node", str(script_dir / "render_content.js"), str(manifest_path)], check=True)
    subprocess.run(["python", str(script_dir / "merge_content.py"), str(work_dir), str(output_path), str(template_path)], check=True)
    if skeleton_pptx.exists():
        skeleton_pptx.unlink()
    print(output_path)


def main():
    if len(sys.argv) < 3:
        print("Usage: python render_report_pptxjs.py content.json output.pptx [template.pptx]")
        sys.exit(1)
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    template = Path(sys.argv[3]) if len(sys.argv) > 3 else skill_dir / "templates" / "吉利模版.pptx"
    build(Path(sys.argv[1]), Path(sys.argv[2]), template)


if __name__ == "__main__":
    main()
