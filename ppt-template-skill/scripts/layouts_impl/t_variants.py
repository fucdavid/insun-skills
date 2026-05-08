from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from ppt_engine.primitives import MUTED, TEXT, rect, rgb, text_fit_size, tx


def _primary(ctx):
    return ctx.token("primary_color", "0B5F68")


def _accent(ctx):
    return ctx.token("accent_color", ctx.token("secondary_color", _primary(ctx)))


def _soft(ctx):
    return ctx.token("bg_light", "F4F7F7")


def _item_text(item):
    if isinstance(item, dict):
        parts = [item.get("desc") or item.get("text") or item.get("detail") or ""]
        extra = item.get("source") or item.get("note") or ""
        if extra:
            parts.append(extra)
        return "\n".join(part for part in parts if part)
    return str(item)


def _item_label(item, fallback):
    if isinstance(item, dict):
        return item.get("label") or item.get("title") or fallback
    return fallback


def _bullets(lines):
    return "\n".join(f"{idx + 1}. {line}" for idx, line in enumerate(lines) if str(line).strip())


def render_cards(slide, body, ctx, start_y):
    primary = _primary(ctx)
    items = (body.get("items") or body.get("cards") or [])[:6]
    if not items:
        items = [{"label": "要点", "desc": body.get("summary") or body.get("statement") or ""}]
    summary = body.get("summary") or body.get("lead") or ""
    if summary:
        rect(slide, 0.72, start_y, 11.9, 0.48, _soft(ctx), _soft(ctx))
        tx(slide, summary, 0.95, start_y + 0.14, 11.4, 0.2, text_fit_size(summary, 11, 8, 58), TEXT, True, "center", ctx=ctx)
        start_y += 0.72
    cols = 3 if len(items) > 2 else max(1, len(items))
    rows = 2 if len(items) > 3 else 1
    gap_x, gap_y = 0.28, 0.28
    card_w = (11.9 - gap_x * (cols - 1)) / cols
    card_h = min(1.55, (4.35 - gap_y * (rows - 1)) / rows)
    for idx, item in enumerate(items):
        col, row = idx % cols, idx // cols
        x = 0.72 + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        rect(slide, x, y, card_w, card_h, "FFFFFF", primary)
        tx(slide, _item_label(item, f"要点{idx + 1}"), x + 0.18, y + 0.18, card_w - 0.36, 0.24, 11, primary, True, ctx=ctx)
        text = _item_text(item)
        tx(slide, text, x + 0.18, y + 0.55, card_w - 0.36, card_h - 0.75, text_fit_size(text, 8.6, 6.8, 42), MUTED, ctx=ctx)


def render_numbered(slide, body, ctx, start_y):
    primary = _primary(ctx)
    items = (body.get("items") or body.get("steps") or [])[:7]
    summary = body.get("summary") or ""
    if summary:
        tx(slide, summary, 0.78, start_y, 11.7, 0.26, text_fit_size(summary, 10.5, 8, 64), MUTED, ctx=ctx)
        start_y += 0.42
    row_h = min(0.72, 4.3 / max(1, len(items)))
    for idx, item in enumerate(items):
        y = start_y + idx * row_h
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.82), Inches(y + 0.1), Inches(0.34), Inches(0.34))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(primary)
        shape.line.fill.background()
        tx(slide, f"{idx + 1:02d}", 0.83, y + 0.18, 0.32, 0.12, 7.5, "FFFFFF", True, "center", ctx=ctx)
        tx(slide, _item_label(item, f"步骤{idx + 1}"), 1.32, y + 0.08, 2.3, 0.22, 10.5, TEXT, True, ctx=ctx)
        tx(slide, _item_text(item), 3.72, y + 0.08, 8.5, 0.34, text_fit_size(_item_text(item), 8.5, 6.6, 58), MUTED, ctx=ctx)


def render_metrics(slide, body, ctx, start_y):
    primary = _primary(ctx)
    metrics = (body.get("metrics") or [])[:5]
    summary = body.get("summary") or ""
    if summary:
        tx(slide, summary, 0.78, start_y, 11.6, 0.28, text_fit_size(summary, 10.5, 8, 64), MUTED, ctx=ctx)
        start_y += 0.52
    cols = min(5, max(1, len(metrics)))
    gap = 0.22
    card_w = (11.9 - gap * (cols - 1)) / cols
    for idx, metric in enumerate(metrics):
        x = 0.72 + idx * (card_w + gap)
        rect(slide, x, start_y, card_w, 1.18, "FFFFFF", primary)
        value = metric.get("value", "") if isinstance(metric, dict) else str(metric)
        label = metric.get("label", f"指标{idx + 1}") if isinstance(metric, dict) else f"指标{idx + 1}"
        desc = metric.get("desc", "") if isinstance(metric, dict) else ""
        tx(slide, value, x + 0.12, start_y + 0.26, card_w - 0.24, 0.26, text_fit_size(value, 17, 10, 12), primary, True, "center", ctx=ctx)
        tx(slide, label, x + 0.12, start_y + 0.68, card_w - 0.24, 0.18, 8.2, TEXT, True, "center", ctx=ctx)
        if desc:
            tx(slide, desc, x + 0.12, start_y + 0.92, card_w - 0.24, 0.14, 6.8, MUTED, False, "center", ctx=ctx)
    details = body.get("details") or body.get("items") or []
    if details:
        tx(slide, _bullets([_item_text(item) for item in details[:6]]), 0.85, start_y + 1.55, 11.5, 1.8, 8.2, MUTED, ctx=ctx)


def render_timeline(slide, body, ctx, start_y):
    primary = _primary(ctx)
    stages = (body.get("stages") or body.get("items") or [])[:6]
    summary = body.get("summary") or ""
    if summary:
        rect(slide, 0.72, start_y, 11.9, 0.46, _soft(ctx), _soft(ctx))
        tx(slide, summary, 0.95, start_y + 0.14, 11.4, 0.18, text_fit_size(summary, 10.5, 8, 60), TEXT, True, "center", ctx=ctx)
        start_y += 0.78
    line_y = start_y + 0.72
    rect(slide, 1.0, line_y, 11.1, 0.03, primary, primary)
    cols = max(1, len(stages))
    gap = 0.22
    card_w = (11.6 - gap * (cols - 1)) / cols
    for idx, stage in enumerate(stages):
        x = 0.88 + idx * (card_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + card_w / 2 - 0.1), Inches(line_y - 0.09), Inches(0.2), Inches(0.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(primary)
        shape.line.fill.background()
        tx(slide, _item_label(stage, f"阶段{idx + 1}"), x, line_y + 0.28, card_w, 0.24, 9, primary, True, "center", ctx=ctx)
        tx(slide, _item_text(stage), x, line_y + 0.68, card_w, 1.35, text_fit_size(_item_text(stage), 7.6, 6.2, 34), MUTED, False, "center", ctx=ctx)


def render_comparison(slide, body, ctx, start_y):
    primary = _primary(ctx)
    columns = body.get("columns") or []
    if not columns:
        columns = [{"title": _item_label(item, f"维度{idx + 1}"), "items": [_item_text(item)]} for idx, item in enumerate((body.get("items") or [])[:4])]
    cols = min(4, max(2, len(columns)))
    gap = 0.22
    col_w = (11.9 - gap * (cols - 1)) / cols
    summary = body.get("summary") or ""
    if summary:
        tx(slide, summary, 0.78, start_y, 11.6, 0.28, text_fit_size(summary, 10.5, 8, 64), MUTED, ctx=ctx)
        start_y += 0.48
    for idx in range(cols):
        col = columns[idx] if idx < len(columns) else {}
        x = 0.72 + idx * (col_w + gap)
        rect(slide, x, start_y, col_w, 0.42, primary, primary)
        tx(slide, col.get("title", f"维度{idx + 1}"), x + 0.12, start_y + 0.13, col_w - 0.24, 0.16, 9, "FFFFFF", True, "center", ctx=ctx)
        rect(slide, x, start_y + 0.52, col_w, 2.9, "FFFFFF", primary)
        values = col.get("items") or col.get("points") or []
        if isinstance(values, str):
            values = [values]
        tx(slide, _bullets(values[:6]), x + 0.18, start_y + 0.75, col_w - 0.36, 2.2, text_fit_size("".join(values), 8, 6.4, 64), MUTED, ctx=ctx)


def render_table(slide, body, ctx, start_y):
    primary = _primary(ctx)
    headers = body.get("headers") or ["事项", "内容", "结果"]
    rows = body.get("rows") or []
    col_count = min(5, len(headers))
    table_w = 11.9
    col_w = table_w / col_count
    rect(slide, 0.72, start_y, table_w, 0.42, primary, primary)
    for idx, header in enumerate(headers[:col_count]):
        tx(slide, str(header), 0.72 + idx * col_w, start_y + 0.13, col_w, 0.15, 8.4, "FFFFFF", True, "center", ctx=ctx)
    row_h = min(0.54, 3.8 / max(1, len(rows)))
    for r, row in enumerate(rows[:7]):
        y = start_y + 0.48 + r * row_h
        rect(slide, 0.72, y, table_w, row_h - 0.04, "FFFFFF", "D9E3E3")
        values = row if isinstance(row, list) else [row.get(h, "") for h in headers]
        for c, value in enumerate(values[:col_count]):
            tx(slide, str(value), 0.78 + c * col_w, y + 0.13, col_w - 0.12, row_h - 0.18, text_fit_size(str(value), 7.2, 5.8, 26), MUTED, ctx=ctx)


def render_statement(slide, body, ctx, start_y):
    primary = _primary(ctx)
    statement = body.get("statement") or body.get("summary") or ""
    supporting = body.get("supporting_text") or "\n".join(body.get("supporting_lines") or [])
    action = body.get("quote") or body.get("action") or ""
    tx(slide, statement, 0.92, start_y + 0.35, 11.4, 0.7, text_fit_size(statement, 24, 14, 42), TEXT, True, "center", ctx=ctx)
    rect(slide, 3.1, start_y + 1.28, 7.2, 0.04, primary, primary)
    if supporting:
        tx(slide, supporting, 1.28, start_y + 1.62, 10.7, 1.25, text_fit_size(supporting, 9.5, 7, 120), MUTED, ctx=ctx)
    if action:
        rect(slide, 1.25, start_y + 3.26, 10.85, 0.52, "111111", "111111")
        tx(slide, action, 1.48, start_y + 3.43, 10.35, 0.15, text_fit_size(action, 12, 8, 48), "FFFFFF", True, "center", ctx=ctx)


def render(slide, body, ctx, start_y=1.35, layout="T4a_cards"):
    variant = body.get("variant") or layout or "T4a_cards"
    if variant == "T4b_numbered_list":
        render_numbered(slide, body, ctx, start_y)
    elif variant == "T4d_metrics":
        render_metrics(slide, body, ctx, start_y)
    elif variant == "T4e_timeline":
        render_timeline(slide, body, ctx, start_y)
    elif variant == "T5_comparison":
        render_comparison(slide, body, ctx, start_y)
    elif variant == "T7_table":
        render_table(slide, body, ctx, start_y)
    elif variant == "T13_statement":
        render_statement(slide, body, ctx, start_y)
    else:
        render_cards(slide, body, ctx, start_y)
