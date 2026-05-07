from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from ppt_engine.primitives import MUTED, TEXT, rect, rgb, tx


def stat(slide, x, y, w, num, label, note, color, ctx):
    rect(slide, x, y, w, 1.32)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(1.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color)
    bar.line.fill.background()
    tx(slide, num, x + 0.2, y + 0.16, w - 0.35, 0.38, 22, color, True, ctx=ctx)
    tx(slide, label, x + 0.22, y + 0.69, w - 0.35, 0.22, 10.8, TEXT, True, ctx=ctx)
    tx(slide, note, x + 0.22, y + 0.98, w - 0.35, 0.22, 8.1, MUTED, ctx=ctx)


def layout_kpi(slide, body, start_y, ctx):
    items = body.get("items", [])
    colors = ctx.palette()
    count = max(1, min(4, len(items)))
    gap = 0.18
    w = (12.1 - gap * (count - 1)) / count
    for i, item in enumerate(items[:4]):
        stat(slide, 0.62 + i * (w + gap), start_y + 0.25, w, item.get("value", ""), item.get("label", ""), item.get("desc", ""), colors[i], ctx)


def layout_grid(slide, body, start_y, ctx):
    items = body.get("items", [])
    colors = ctx.palette()
    for i, item in enumerate(items[:6]):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.1
        y = start_y + row * 1.22
        rect(slide, x, y, 5.35, 0.92)
        tx(slide, f"{i + 1:02d}", x + 0.25, y + 0.2, 0.45, 0.18, 9.5, colors[i % len(colors)], True, ctx=ctx)
        tx(slide, item.get("label", ""), x + 0.9, y + 0.18, 2.6, 0.18, 10.5, TEXT, True, ctx=ctx)
        tx(slide, item.get("desc", ""), x + 0.9, y + 0.52, 3.95, 0.32, 8.4, MUTED, ctx=ctx)


def layout_timeline(slide, body, start_y, ctx):
    items = body.get("items", [])
    primary = ctx.token("primary_color", "093BAA")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(start_y + 0.72), Inches(11.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(ctx.token("bg_light", "C8E0FB"))
    line.line.fill.background()
    step_w = 10.8 / max(1, len(items))
    for i, item in enumerate(items[:6]):
        cx = 1.35 + i * step_w + step_w / 2
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.08), Inches(start_y + 0.62), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(primary)
        dot.line.fill.background()
        tx(slide, item.get("label", ""), cx - step_w / 2 + 0.05, start_y + 0.93, step_w - 0.1, 0.2, 10.2, TEXT, True, "center", ctx=ctx)
        tx(slide, item.get("desc", ""), cx - step_w / 2 + 0.1, 4.28, step_w - 0.2, 0.6, 8.2, MUTED, False, "center", ctx=ctx)


def layout_big(slide, body, start_y, ctx):
    item = (body.get("items") or [{}])[0]
    primary = ctx.token("primary_color", "093BAA")
    tx(slide, item.get("value", ""), 1.45, start_y + 0.3, 10.2, 0.75, 40, primary, True, "center", "number", ctx)
    tx(slide, item.get("label", ""), 1.45, start_y + 1.3, 10.2, 0.35, 18, TEXT, True, "center", ctx=ctx)
    tx(slide, item.get("desc", ""), 1.75, start_y + 2.0, 9.6, 0.6, 12, MUTED, False, "center", ctx=ctx)


def layout_insight(slide, body, start_y, ctx):
    primary = ctx.token("primary_color", "093BAA")
    insight = body.get("insight") or body.get("headline") or ""
    quote = body.get("quote") or ""
    rect(slide, 0.72, start_y + 0.05, 4.6, 0.92)
    tx(slide, "核心观点", 0.96, start_y + 0.25, 0.85, 0.18, 9, primary, True, ctx=ctx)
    tx(slide, insight, 1.82, start_y + 0.23, 2.9, 0.55, 10.2, TEXT, True, ctx=ctx)
    tx(slide, body.get("path_title") or "关键演进路径", 0.72, start_y + 1.47, 4.1, 0.22, 11.2, TEXT, True, ctx=ctx)
    for i, item in enumerate((body.get("path") or body.get("evolution") or [])[:5]):
        py = start_y + 1.88 + i * 0.43
        tx(slide, item.get("label", ""), 1.58, py + 0.01, 1.25, 0.17, 8.7, TEXT, True, ctx=ctx)
        tx(slide, item.get("desc", ""), 2.78, py + 0.01, 2.25, 0.17, 6.8, MUTED, ctx=ctx)
    tx(slide, body.get("data_title") or "关键数据", 5.7, start_y + 0.05, 4.0, 0.22, 11.2, TEXT, True, ctx=ctx)
    for i, item in enumerate((body.get("data") or [])[:4]):
        x = 5.72 + (i % 2) * 2.65
        y = start_y + 0.52 + (i // 2) * 0.95
        rect(slide, x, y, 2.35, 0.76)
        tx(slide, item.get("label", ""), x + 0.2, y + 0.12, 1.9, 0.18, 7.8, TEXT, True, ctx=ctx)
        tx(slide, item.get("desc", ""), x + 0.2, y + 0.39, 1.9, 0.24, 6.3, MUTED, ctx=ctx)
    if quote:
        rect(slide, 5.72, start_y + 2.77, 5.25, 0.98)
        tx(slide, quote, 6.05, start_y + 3.03, 4.55, 0.34, 13.5, primary, True, "center", ctx=ctx)


def render(slide, body, ctx, start_y=1.35, layout="grid"):
    if layout == "kpi":
        layout_kpi(slide, body, start_y, ctx)
    elif layout == "timeline":
        layout_timeline(slide, body, start_y, ctx)
    elif layout == "big_stat":
        layout_big(slide, body, max(1.75, start_y), ctx)
    elif layout == "insight_onepage":
        layout_insight(slide, body, start_y, ctx)
    else:
        layout_grid(slide, body, start_y, ctx)

