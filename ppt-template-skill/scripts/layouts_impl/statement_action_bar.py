from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from ppt_engine.primitives import MUTED, TEXT, rich_tx, rgb, text_fit_size, tx


LAYOUT_NAME = "statement_action_bar"
OWNS_TITLE = True


def render(slide, body, ctx, start_y=1.35, layout=None):
    statement = body.get("statement") or body.get("headline") or body.get("title") or ""
    emphasis = body.get("emphasis") or ""
    action = body.get("action") or body.get("conclusion") or body.get("footer") or ""
    action_highlight = body.get("action_highlight") or ""
    primary = ctx.token("primary_color", "093BAA")
    text_color = ctx.token("text_dark", TEXT)
    muted_color = ctx.token("text_mid", MUTED)
    highlight_color = body.get("highlight_color") or ctx.token("highlight_color", "FFF200")

    anchor_x = 1.14
    anchor_y = 2.81
    text_x = 2.08
    prefix_y = 2.88
    emphasis_y = 3.21
    text_w = 9.25

    icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(anchor_x + 0.07), Inches(anchor_y + 0.08), Inches(0.34), Inches(0.34))
    icon.fill.solid()
    icon.fill.fore_color.rgb = rgb("6E747A")
    icon.line.fill.background()
    lens = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(anchor_x + 0.18), Inches(anchor_y + 0.18), Inches(0.105), Inches(0.105))
    lens.fill.background()
    lens.line.color.rgb = rgb("FFFFFF")
    lens.line.width = Pt(1.0)
    handle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(anchor_x + 0.275), Inches(anchor_y + 0.285), Inches(0.065), Inches(0.015))
    handle.rotation = 45
    handle.fill.solid()
    handle.fill.fore_color.rgb = rgb("FFFFFF")
    handle.line.fill.background()

    for dot_i in range(3):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(anchor_x + 0.21), Inches(anchor_y + 0.60 + dot_i * 0.25), Inches(0.065), Inches(0.065))
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(muted_color)
        dot.fill.transparency = 10
        dot.line.fill.background()

    if emphasis and emphasis in statement:
        prefix, suffix = statement.split(emphasis, 1)
        prefix_text = f"{prefix}{suffix}".strip()
        tx(slide, prefix_text, text_x, prefix_y, text_w, 0.32, text_fit_size(prefix_text, 14, 10, 36), text_color, False, "left", ctx=ctx)
        tx(slide, emphasis, text_x, emphasis_y, text_w, 0.55, text_fit_size(emphasis, 27, 18, 24), text_color, True, "left", "title", ctx)
    else:
        tx(slide, statement, text_x, emphasis_y, text_w, 0.65, text_fit_size(statement, 25, 17, 30), text_color, True, "left", "title", ctx)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM, Inches(6.16), Inches(4.02), Inches(5.18), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(primary)
    line.line.fill.background()

    bar_y = 4.15
    action_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.07), Inches(bar_y), Inches(9.32), Inches(0.56))
    action_bar.fill.solid()
    action_bar.fill.fore_color.rgb = rgb("000000")
    action_bar.line.fill.background()

    if action_highlight and action_highlight in action:
        before, after = action.split(action_highlight, 1)
        rich_tx(
            slide,
            [
                {"text": before, "size": 13, "bold": True, "color": "FFFFFF"},
                {"text": action_highlight, "size": 13, "bold": True, "color": highlight_color},
                {"text": after, "size": 13, "bold": True, "color": "FFFFFF"},
            ],
            2.18,
            bar_y + 0.16,
            9.05,
            0.18,
            13,
            "center",
            "body",
            ctx,
        )
    else:
        tx(slide, action, 2.18, bar_y + 0.16, 9.05, 0.18, 13, "FFFFFF", True, "center", ctx=ctx)

    if body.get("note"):
        tx(slide, body.get("note"), 2.0, bar_y + 0.75, 9.0, 0.18, 8, muted_color, False, "center", ctx=ctx)

