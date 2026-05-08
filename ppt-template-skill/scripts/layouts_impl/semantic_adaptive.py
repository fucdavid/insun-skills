from __future__ import annotations

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from ppt_engine.primitives import MUTED, TEXT, rect, rgb, text_fit_size, tx


def _primary(ctx):
    return ctx.token("primary_color", "861B2F")


def _accent(ctx):
    return ctx.token("accent_color", ctx.token("primary_color", "861B2F"))


def _soft(ctx):
    return ctx.token("bg_light", "F5F2EE")


def _line(slide, x, y, w, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.025))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def _dot(slide, x, y, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.13), Inches(0.13))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def render_statement(slide, body, ctx, start_y):
    primary = _primary(ctx)
    statement = body.get("statement") or body.get("summary") or ""
    emphasis = body.get("emphasis") or statement
    action = body.get("action") or body.get("action_highlight") or ""
    _dot(slide, 0.86, start_y + 0.52, primary)
    for i in range(3):
        _dot(slide, 0.9, start_y + 0.84 + i * 0.2, "777777")
    tx(slide, statement, 1.25, start_y + 0.34, 10.8, 0.55, text_fit_size(statement, 20, 13, 42), TEXT, False, ctx=ctx)
    tx(slide, emphasis, 1.25, start_y + 0.9, 10.8, 0.55, text_fit_size(emphasis, 25, 15, 30), primary, True, ctx=ctx)
    _line(slide, 4.85, start_y + 1.58, 5.2, primary)
    rect(slide, 1.25, start_y + 1.82, 10.8, 0.62, "111111", "111111")
    tx(slide, action, 1.48, start_y + 2.01, 10.3, 0.23, text_fit_size(action, 14, 9.5, 42), "FFFFFF", True, "center", ctx=ctx)


def render_strategy_timeline(slide, body, ctx, start_y):
    primary = _primary(ctx)
    summary = body.get("summary", "")
    rect(slide, 0.72, start_y + 0.02, 11.8, 0.54, _soft(ctx), primary)
    tx(slide, summary, 0.95, start_y + 0.18, 11.35, 0.18, text_fit_size(summary, 11, 8, 58), TEXT, True, "center", ctx=ctx)
    tx(slide, body.get("strategy_title", "Strategy"), 0.72, start_y + 0.78, 3.1, 0.24, 12, primary, True, ctx=ctx)
    tx(slide, body.get("core_point", ""), 3.05, start_y + 0.78, 8.8, 0.24, 11, TEXT, True, ctx=ctx)
    stages = (body.get("stages") or [])[:3]
    gap = 0.28
    card_w = (11.8 - gap * 2) / 3
    for i in range(3):
        stage = stages[i] if i < len(stages) else {}
        x = 0.72 + i * (card_w + gap)
        y = start_y + 1.3
        rect(slide, x, y, card_w, 2.15, "FFFFFF", primary)
        tx(slide, stage.get("period", f"Stage {i+1}"), x + 0.18, y + 0.18, card_w - 0.36, 0.22, 12, primary, True, "center", ctx=ctx)
        tx(slide, stage.get("theme", ""), x + 0.18, y + 0.58, card_w - 0.36, 0.28, 13, TEXT, True, "center", ctx=ctx)
        tx(slide, stage.get("goal", ""), x + 0.28, y + 1.05, card_w - 0.56, 0.68, text_fit_size(stage.get("goal", ""), 9.5, 7, 42), MUTED, ctx=ctx)
    tx(slide, body.get("explanation", ""), 0.86, start_y + 3.72, 11.45, 0.24, 10.2, MUTED, False, "center", ctx=ctx)


def render_phase_plan(slide, body, ctx, start_y):
    primary = _primary(ctx)
    tx(slide, body.get("summary", ""), 0.72, start_y + 0.02, 11.8, 0.34, 10, MUTED, ctx=ctx)
    stages = (body.get("stages") or [])[:3]
    for i in range(3):
        stage = stages[i] if i < len(stages) else {}
        y = start_y + 0.55 + i * 1.24
        rect(slide, 0.82, y, 11.5, 0.95, "FFFFFF", primary)
        tx(slide, stage.get("phase", f"Step {i+1}"), 1.02, y + 0.16, 0.9, 0.2, 10, primary, True, ctx=ctx)
        tx(slide, stage.get("event", ""), 2.05, y + 0.13, 4.5, 0.22, 10, TEXT, True, ctx=ctx)
        tx(slide, stage.get("strategy", ""), 2.05, y + 0.48, 4.5, 0.22, 8.2, MUTED, ctx=ctx)
        tx(slide, stage.get("keyword", ""), 6.85, y + 0.14, 2.25, 0.22, 9.5, primary, True, ctx=ctx)
        tx(slide, stage.get("detail", ""), 9.15, y + 0.14, 2.75, 0.45, 8, MUTED, ctx=ctx)
    tx(slide, body.get("fee", ""), 0.92, start_y + 4.32, 11.2, 0.24, 9.2, TEXT, True, "center", ctx=ctx)


def render_risk_matrix(slide, body, ctx, start_y):
    primary = _primary(ctx)
    tx(slide, body.get("summary", ""), 0.72, start_y + 0.02, 11.8, 0.32, 9.2, MUTED, ctx=ctx)
    labels = body.get("column_labels") or ["风险预判", "应对策略", "行动抓手", "监测机制"]
    subjects = body.get("risk_subjects") or []
    cols = [body.get("risk_forecasts") or [], body.get("response_strategies") or [], body.get("actions") or [], body.get("monitoring_guidance") or []]
    x0, y0 = 0.68, start_y + 0.52
    row_h = 0.78
    widths = [1.6, 2.55, 2.55, 2.4, 2.1]
    tx(slide, body.get("risk_subject_label", "主题"), x0, y0, widths[0], 0.24, 9.5, primary, True, "center", ctx=ctx)
    x = x0 + widths[0]
    for i, label in enumerate(labels[:4]):
        tx(slide, label, x, y0, widths[i + 1], 0.24, 9.5, primary, True, "center", ctx=ctx)
        x += widths[i + 1]
    for r in range(4):
        y = y0 + 0.42 + r * row_h
        rect(slide, x0, y, sum(widths), row_h - 0.08, "FFFFFF", primary)
        tx(slide, subjects[r] if r < len(subjects) else f"Topic {r+1}", x0 + 0.12, y + 0.16, widths[0] - 0.24, 0.22, 8.8, TEXT, True, "center", ctx=ctx)
        x = x0 + widths[0]
        for c in range(4):
            values = cols[c]
            text = values[r] if r < len(values) else ""
            tx(slide, text, x + 0.12, y + 0.12, widths[c + 1] - 0.24, 0.34, text_fit_size(text, 7.7, 6.2, 28), MUTED, ctx=ctx)
            x += widths[c + 1]


def render_metric_summary(slide, body, ctx, start_y):
    primary = _primary(ctx)
    tx(slide, body.get("summary", ""), 0.72, start_y + 0.02, 11.8, 0.32, 9.5, MUTED, ctx=ctx)
    metrics = (body.get("metrics") or [])[:3]
    for i in range(3):
        metric = metrics[i] if i < len(metrics) else {}
        x = 0.82 + i * 3.9
        rect(slide, x, start_y + 0.5, 3.45, 0.9, "FFFFFF", primary)
        tx(slide, metric.get("value", ""), x + 0.18, start_y + 0.68, 3.05, 0.24, text_fit_size(metric.get("value", ""), 16, 10, 12), primary, True, "center", ctx=ctx)
        tx(slide, metric.get("label", ""), x + 0.18, start_y + 1.03, 3.05, 0.18, 8.5, TEXT, True, "center", ctx=ctx)
    topics = (body.get("topics") or [])[:3]
    for i in range(3):
        topic = topics[i] if i < len(topics) else {}
        x = 0.82 + i * 3.9
        rect(slide, x, start_y + 1.72, 3.45, 2.15, "FFFFFF", primary)
        tx(slide, topic.get("title", f"Topic {i+1}"), x + 0.18, start_y + 1.92, 3.05, 0.2, 10, TEXT, True, "center", ctx=ctx)
        events = topic.get("events") or []
        lines = []
        for event in events[:4]:
            if isinstance(event, dict):
                lines.append((event.get("date", "") + " " + event.get("result", "")).strip())
            else:
                lines.append(str(event))
        tx(slide, "\n".join(lines), x + 0.22, start_y + 2.32, 3.0, 0.95, 7.4, MUTED, ctx=ctx)
    tx(slide, body.get("issue", ""), 0.92, start_y + 4.22, 11.2, 0.24, 9.2, TEXT, True, "center", ctx=ctx)


def render(slide, body, ctx, start_y=1.35, layout="grid"):
    if layout == "statement_action_bar":
        render_statement(slide, body, ctx, start_y)
    elif layout == "style2_overall_strategy_timeline":
        render_strategy_timeline(slide, body, ctx, start_y)
    elif layout == "sample3_ota_timeline":
        render_phase_plan(slide, body, ctx, start_y)
    elif layout == "style2_risk_response_matrix":
        render_risk_matrix(slide, body, ctx, start_y)
    elif layout == "sample3_koc_monthly_summary":
        render_metric_summary(slide, body, ctx, start_y)
    else:
        render_statement(slide, body, ctx, start_y)
