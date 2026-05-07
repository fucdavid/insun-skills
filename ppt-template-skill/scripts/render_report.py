#!/usr/bin/env python3
from __future__ import annotations

import json
import random
import sys
from pathlib import Path

from pptx import Presentation

from ppt_engine.content import (
    ensure_toc_slide,
    render_content_slide,
    render_toc_slide,
    render_transition_slide,
    replace_cover_text,
)
from ppt_engine.context import RenderContext
from ppt_engine.layout_registry import LayoutRegistry
from ppt_engine.template import (
    collect_slide_fonts,
    delete_slide,
    duplicate_slide,
    load_template_tokens,
    move_slide,
    prepare_content_template,
    resolve_template_roles,
)


def build(content_path: Path, output_path: Path, template_path: Path):
    content = json.loads(content_path.read_text(encoding="utf-8-sig"))
    prs = Presentation(str(template_path))
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    registry = LayoutRegistry(skill_dir)
    ctx = RenderContext(skill_dir=skill_dir, template_path=template_path, tokens=load_template_tokens(template_path))

    cover_idx, toc_template_idx, transition_indices, content_indices, end_idx = resolve_template_roles(prs, template_path)
    inferred_fonts = collect_slide_fonts([prs.slides[idx] for idx in content_indices if idx < len(prs.slides)])
    if inferred_fonts and not ctx.tokens.get("font_body"):
        ctx.tokens["font_body"] = inferred_fonts.most_common(1)[0][0]
    if inferred_fonts and not ctx.tokens.get("font_title"):
        ctx.tokens["font_title"] = ctx.tokens.get("font_body")

    desired_slides = ensure_toc_slide(content.get("slides") or [], toc_template_idx is not None)
    body_slides = [s for s in desired_slides if s.get("type") in ("toc", "transition", "chapter", "content")]
    toc_source = prs.slides[toc_template_idx] if toc_template_idx is not None else prs.slides[content_indices[0]]
    transition_sources = [prs.slides[idx] for idx in transition_indices] or [prs.slides[content_indices[0]]]
    content_sources = [prs.slides[idx] for idx in content_indices]
    end_source = prs.slides[end_idx]
    cover_slide_id = prs.slides[cover_idx].slide_id
    end_slide_id = end_source.slide_id

    style_selection = content.get("style_selection", "rotate")
    rng = random.Random(content.get("style_seed", 20260430))
    generated = []
    content_style_cursor = 0
    transition_style_cursor = 0
    for slide_data in body_slides:
        slide_type = slide_data.get("type")
        if slide_type == "toc":
            generated.append((duplicate_slide(prs, toc_source), slide_data))
        elif slide_type in ("transition", "chapter"):
            source = transition_sources[transition_style_cursor % len(transition_sources)]
            transition_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))
        else:
            if style_selection == "random":
                source = rng.choice(content_sources)
            else:
                source = content_sources[content_style_cursor % len(content_sources)]
                content_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))

    keep_ids = {cover_slide_id, end_slide_id}
    keep_ids.update(slide.slide_id for slide, _ in generated)
    for idx in range(len(prs.slides) - 1, -1, -1):
        if prs.slides[idx].slide_id not in keep_ids:
            delete_slide(prs, idx)

    target_order = [cover_slide_id] + [slide.slide_id for slide, _ in generated] + [end_slide_id]
    for target_idx, slide_id in enumerate(target_order):
        current_idx = next(i for i, slide in enumerate(prs.slides) if slide.slide_id == slide_id)
        if current_idx != target_idx:
            move_slide(prs, current_idx, target_idx)

    slides = list(prs.slides)
    template_title_bottom = {}
    for slide, slide_data in generated:
        if slide_data.get("type") in ("toc", "transition", "chapter"):
            continue
        template_title_bottom[slide.slide_id] = prepare_content_template(slide, slide_data, registry)

    cover_data = next((s for s in desired_slides if s.get("type") == "cover"), {})
    replace_cover_text(slides[0], cover_data, content)

    source = content.get("source", "")
    for i, slide_data in enumerate(body_slides, start=1):
        slide = slides[i]
        if slide_data.get("type") == "toc":
            render_toc_slide(slide, slide_data, i + 1, source, ctx)
        elif slide_data.get("type") in ("transition", "chapter"):
            render_transition_slide(slide, slide_data, i + 1, source, ctx)
        else:
            render_content_slide(
                slide,
                slide_data,
                i + 1,
                source,
                registry,
                ctx,
                draw_title=template_title_bottom.get(slide.slide_id) is None,
                content_top=template_title_bottom.get(slide.slide_id),
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(output_path)


def main():
    if len(sys.argv) < 3:
        print("Usage: python render_report.py content.json output.pptx [template.pptx]")
        sys.exit(1)
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    template = Path(sys.argv[3]) if len(sys.argv) > 3 else skill_dir / "templates" / "吉利模版.pptx"
    build(Path(sys.argv[1]), Path(sys.argv[2]), template)


if __name__ == "__main__":
    main()
