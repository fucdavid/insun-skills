from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
TEXT = "000000"
MUTED = "333333"
PANEL = "FFFFFF"
LIGHT_BLUE = "C8E0FB"


def emu_to_in(value):
    return float(value) / EMU_PER_INCH


def rgb(hex_color: str) -> RGBColor:
    hex_color = str(hex_color).replace("#", "")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def set_run_font(run, font_name: str):
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", font_name)
        node.set("charset", "-122")


def tx(slide, text, x, y, w, h, size=16, color=TEXT, bold=False, align="left", font_role="body", ctx=None):
    font_name = ctx.font_for(font_role) if ctx else "Microsoft YaHei"
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = str(text or "")
    set_run_font(run, font_name)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def rich_tx(slide, segments, x, y, w, h, size=16, align="left", font_role="body", ctx=None):
    font_name = ctx.font_for(font_role) if ctx else "Microsoft YaHei"
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    for segment in segments:
        run = p.add_run()
        run.text = str(segment.get("text", ""))
        set_run_font(run, font_name)
        run.font.size = Pt(segment.get("size", size))
        run.font.bold = bool(segment.get("bold", False))
        run.font.color.rgb = rgb(segment.get("color", TEXT))
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=LIGHT_BLUE, alpha=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = alpha
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.6)
    return shape


def title(slide, heading, subtitle=None, ctx=None):
    tx(slide, heading, 0.55, 0.34, 10.8, 0.38, 22, TEXT, True, font_role="title", ctx=ctx)
    if subtitle:
        tx(slide, subtitle, 0.56, 0.86, 11.2, 0.24, 10.2, MUTED, ctx=ctx)


def footer(slide, page, source, ctx=None):
    if source:
        tx(slide, source, 0.55, 7.08, 9.4, 0.16, 7, MUTED, ctx=ctx)
    tx(slide, f"{page:02d}", 12.24, 7.04, 0.38, 0.16, 7.5, MUTED, True, "right", ctx=ctx)


def text_fit_size(text, base, minimum, soft_limit):
    length = len(str(text or ""))
    if length <= soft_limit:
        return base
    return max(minimum, base - (length - soft_limit) * 0.42)


def add_image_or_slot(slide, image_path, x, y, w, h, label="素材图", ctx=None):
    if image_path:
        p = Path(str(image_path))
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
            return
    slot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    slot.fill.solid()
    slot.fill.fore_color.rgb = rgb("EEF2F7")
    slot.line.color.rgb = rgb("CCD6E6")
    slot.line.width = Pt(0.6)
    tx(slide, label, x + 0.18, y + h / 2 - 0.08, max(0.2, w - 0.36), 0.16, 8, MUTED, True, "center", ctx=ctx)

