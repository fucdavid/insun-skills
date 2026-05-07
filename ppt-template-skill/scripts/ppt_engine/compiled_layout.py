from __future__ import annotations

import json
from io import BytesIO
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from .primitives import add_image_or_slot, emu_to_in
from .template import iter_text_shapes, replace_text_preserve_shape


_CACHE: dict[Path, Presentation] = {}


def _load_presentation(path: Path) -> Presentation:
    resolved = path.resolve()
    if resolved not in _CACHE:
        _CACHE[resolved] = Presentation(str(resolved))
    return _CACHE[resolved]


def _replace_rel_ids(element, rel_id_map: dict[str, str]):
    for node in element.iter():
        for attr_name, attr_value in list(node.attrib.items()):
            if attr_value in rel_id_map:
                node.set(attr_name, rel_id_map[attr_value])


def _strip_placeholder_binding(element):
    # Placeholder shapes can be re-bound to the target template layout by
    # PowerPoint, changing position/orientation. Compiled layouts need source
    # geometry to stay literal, so convert placeholders to ordinary shapes.
    for ph in list(element.iter(qn("p:ph"))):
        parent = ph.getparent()
        if parent is not None:
            parent.remove(ph)


def _normalize_sp_tree_base(sp_tree):
    base_tags = [qn("p:nvGrpSpPr"), qn("p:grpSpPr")]
    for target_index, tag in enumerate(base_tags):
        child = next((candidate for candidate in list(sp_tree) if candidate.tag == tag), None)
        if child is None:
            continue
        current_parent = child.getparent()
        if current_parent is not None:
            current_parent.remove(child)
        sp_tree.insert(target_index, child)


def _insert_shape_before_ext_list(sp_tree, element):
    ext_lst = sp_tree.find(qn("p:extLst"))
    if ext_lst is None:
        sp_tree.append(element)
        return
    sp_tree.insert(sp_tree.index(ext_lst), element)


def _ensure_explicit_geometry(element, shape):
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except Exception:
        return
    if None in (left, top, width, height):
        return

    sp_pr = element.find(qn("p:spPr"))
    if sp_pr is None:
        sp_pr = OxmlElement("p:spPr")
        element.insert(1, sp_pr)

    xfrm = sp_pr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = OxmlElement("a:xfrm")
        sp_pr.insert(0, xfrm)

    off = xfrm.find(qn("a:off"))
    if off is None:
        off = OxmlElement("a:off")
        xfrm.insert(0, off)
    off.set("x", str(int(left)))
    off.set("y", str(int(top)))

    ext = xfrm.find(qn("a:ext"))
    if ext is None:
        ext = OxmlElement("a:ext")
        xfrm.append(ext)
    ext.set("cx", str(int(width)))
    ext.set("cy", str(int(height)))


def _copy_source_shapes(target_slide, source_slide):
    rel_id_map = {}
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype or rel.reltype.endswith("/tags"):
            continue
        if rel.reltype.endswith("/image") and not rel.is_external:
            try:
                _, new_rid = target_slide.part.get_or_add_image_part(BytesIO(rel.target_part.blob))
                rel_id_map[rel.rId] = new_rid
                continue
            except Exception:
                pass
        if rel.is_external:
            new_rid = target_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = target_slide.part.relate_to(rel.target_part, rel.reltype)
        rel_id_map[rel.rId] = new_rid

    for shape in list(target_slide.shapes):
        shape.element.getparent().remove(shape.element)
    _normalize_sp_tree_base(target_slide.shapes._spTree)

    for shape in source_slide.shapes:
        element = deepcopy(shape.element)
        _replace_rel_ids(element, rel_id_map)
        _ensure_explicit_geometry(element, shape)
        _strip_placeholder_binding(element)
        _insert_shape_before_ext_list(target_slide.shapes._spTree, element)


def _get_path_value(data, path: str):
    current = data
    for part in path.split("."):
        if "[" in part and part.endswith("]"):
            name, index_text = part[:-1].split("[", 1)
            current = current.get(name, []) if isinstance(current, dict) else []
            try:
                current = current[int(index_text)]
            except (IndexError, ValueError, TypeError):
                return ""
            continue
        if isinstance(current, dict):
            current = current.get(part, "")
        else:
            return ""
    return current


def _format_value(value):
    if value is None:
        return ""
    if isinstance(value, list):
        return "\n".join(_format_value(item) for item in value)
    if isinstance(value, dict):
        title = value.get("title") or value.get("label") or value.get("date") or ""
        desc = value.get("desc") or value.get("result") or value.get("summary") or ""
        if title and desc:
            return f"{title}\n{desc}"
        return str(title or desc)
    return str(value)


def _iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        try:
            is_group = getattr(shape, "shape_type", None) and str(shape.shape_type).startswith("GROUP")
        except NotImplementedError:
            is_group = False
        if is_group:
            yield from _iter_all_shapes(shape.shapes)


def _replace_text_frame_text(text_frame, text):
    text = "" if text is None else str(text)
    first_run = None
    first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            first_run = paragraph.runs[0]
            break
    if first_run is None:
        if first_paragraph is None:
            text_frame.text = text
            return
        first_run = first_paragraph.add_run()

    first_run.text = text
    first_run_element = first_run._r
    seen_first = False
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run._r is first_run_element and not seen_first:
                seen_first = True
                continue
            run.text = ""


def _replace_table(table_shape, value):
    if not getattr(table_shape, "has_table", False):
        return
    rows = value if isinstance(value, list) else []
    table = table_shape.table
    for r_idx, row in enumerate(table.rows):
        row_values = rows[r_idx] if r_idx < len(rows) else []
        if isinstance(row_values, dict):
            row_values = list(row_values.values())
        if not isinstance(row_values, list):
            row_values = [row_values]
        for c_idx, cell in enumerate(row.cells):
            cell_value = row_values[c_idx] if c_idx < len(row_values) else ""
            _replace_text_frame_text(cell.text_frame, _format_value(cell_value))


def _drop_shape_image_rels(slide, shape):
    for blip in shape.element.iter(qn("a:blip")):
        rel_id = blip.get(qn("r:embed"))
        if not rel_id:
            continue
        try:
            slide.part.drop_rel(rel_id)
        except KeyError:
            pass


def _insert_new_shapes_at_original_z_order(slide, original_shape, add_fn):
    parent = original_shape.element.getparent()
    if parent is None:
        add_fn()
        return
    root = slide.shapes._spTree
    insert_parent = parent
    insert_index = parent.index(original_shape.element)

    # If the source picture lives inside a group, the replacement shape created
    # by python-pptx must stay at slide root level. Moving a root-level shape
    # into a group can corrupt the grouped shape XML and change z-order. In this
    # case insert the replacement at the group's root-level z position instead.
    ancestor = original_shape.element
    while ancestor.getparent() is not None and ancestor.getparent() is not root:
        ancestor = ancestor.getparent()
    if parent is not root and ancestor.getparent() is root:
        insert_parent = root
        insert_index = root.index(ancestor)

    before_child_ids = {id(child) for child in insert_parent}

    parent.remove(original_shape.element)
    add_fn()

    added_elements = [child for child in list(insert_parent) if id(child) not in before_child_ids]
    for element in added_elements:
        insert_parent.remove(element)
    for offset, element in enumerate(added_elements):
        insert_parent.insert(insert_index + offset, element)


def render_compiled_layout(slide, body, ctx, layout_name: str):
    spec_path = ctx.skill_dir / "layouts" / "specs" / f"{layout_name}.json"
    spec = json.loads(spec_path.read_text(encoding="utf-8-sig"))
    source_path = ctx.skill_dir / spec["compiled_pptx"]
    source_prs = _load_presentation(source_path)
    source_slide = source_prs.slides[int(spec.get("source_slide", 1)) - 1]

    _copy_source_shapes(slide, source_slide)

    text_shape_list = list(iter_text_shapes(slide.shapes))
    by_name = {getattr(shape, "name", ""): shape for shape in text_shape_list}
    for field in spec.get("fields", []):
        shape = None
        if field.get("match_text") is not None:
            shape = next((candidate for candidate in text_shape_list if (candidate.text or "").strip() == field.get("match_text")), None)
        if shape is None:
            shape = by_name.get(field.get("shape", ""))
        if shape is None:
            continue
        value = _get_path_value(body, field.get("path", ""))
        text = _format_value(value)
        if not text and field.get("clear_if_empty", True):
            replace_text_preserve_shape(shape, "")
            continue
        if text:
            replace_text_preserve_shape(shape, text)

    for shape_name in spec.get("clear_shapes", []):
        shape = by_name.get(shape_name)
        if shape is not None:
            replace_text_preserve_shape(shape, "")

    all_shapes = list(_iter_all_shapes(slide.shapes))
    by_shape_name = {getattr(shape, "name", ""): shape for shape in all_shapes}
    for table_field in spec.get("table_fields", []):
        shape = by_shape_name.get(table_field.get("shape", ""))
        if shape is None:
            continue
        _replace_table(shape, _get_path_value(body, table_field.get("path", "")))

    for image_field in spec.get("image_fields", []):
        shape = by_shape_name.get(image_field.get("shape", ""))
        if shape is None:
            continue
        image_path = _get_path_value(body, image_field.get("path", ""))
        if not image_path or not Path(str(image_path)).exists():
            continue
        x = emu_to_in(shape.left)
        y = emu_to_in(shape.top)
        w = emu_to_in(shape.width)
        h = emu_to_in(shape.height)
        _drop_shape_image_rels(slide, shape)
        _insert_new_shapes_at_original_z_order(
            slide,
            shape,
            lambda: add_image_or_slot(slide, image_path, x, y, w, h, image_field.get("label", "素材图"), ctx),
        )
