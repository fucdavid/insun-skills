#!/usr/bin/env python3
"""
Generate a report from the built-in Geely template.

Use this path for templates shaped like:
  slide1 = cover, slide2 = blank/content style, slide3 = closing/quote page

Usage:
  python scripts/render_geely_report.py content.json output.pptx templates/吉利模版.pptx

content.json shape:
{
  "title": "小鹏汽车2025年发展总结",
  "date": "2026.04",
  "slides": [
    {"type": "cover", "title": "...", "date": "..."},
    {
      "type": "content",
      "page_title": "核心结论",
      "subtitle": "可选副标题",
      "body": {
        "layout": "auto",
        "items": [
          {"label": "全年交付", "value": "42.94万", "desc": "同比增长125.9%"},
          {"label": "产品周期", "desc": "MONA M03、P7+、G7共同支撑"}
        ]
      }
    },
    {"type": "end", "title": "谢谢", "date": "2026.04"}
  ]
}
"""

from __future__ import annotations

import json
import random
import re
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


BLUE = "093BAA"
GREEN = "10A875"
CYAN = "2DBCEB"
ORANGE = "F28C28"
LIGHT_BLUE = "C8E0FB"
TEXT = "000000"
MUTED = "333333"
PANEL = "FFFFFF"
FONT = "汉仪雅酷黑 75W"
EMU_PER_INCH = 914400


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def emu_to_in(value):
    return float(value) / EMU_PER_INCH


def tx(slide, text, x, y, w, h, size=16, color=TEXT, bold=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)

    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", FONT)
        node.set("charset", "-122")
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=LIGHT_BLUE, alpha=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = alpha
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.6)
    return shape


def title(slide, heading, subtitle=None):
    tx(slide, heading, 0.55, 0.34, 10.8, 0.38, 22, TEXT, True)
    if subtitle:
        tx(slide, subtitle, 0.56, 0.86, 11.2, 0.24, 10.2, MUTED)


def footer(slide, page, source):
    if source:
        tx(slide, source, 0.55, 7.08, 9.4, 0.16, 7, MUTED)
    tx(slide, f"{page:02d}", 12.24, 7.04, 0.38, 0.16, 7.5, MUTED, True, "right")


def replace_text_preserve_shape(shape, text):
    tf = shape.text_frame
    first_run = None
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if first_run is None:
                first_run = run
            else:
                run.text = ""
    if first_run is not None:
        first_run.text = text
    else:
        shape.text = text


def text_shapes(slide):
    return [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    ]


def all_text_frame_shapes(slide):
    shapes = []

    def walk(shape_collection):
        for shape in shape_collection:
            if getattr(shape, "has_text_frame", False):
                shapes.append(shape)
            if hasattr(shape, "shapes"):
                walk(shape.shapes)

    walk(slide.shapes)
    return shapes


def iter_text_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            yield shape
        if hasattr(shape, "shapes"):
            yield from iter_text_shapes(shape.shapes)


def is_placeholder(shape):
    return shape.element.xpath(".//p:ph") != []


def clear_slide_shapes(slide, preserve_pictures=False, preserve_decorations=False):
    for shape in list(slide.shapes):
        if preserve_pictures and shape.element.tag.endswith('}pic'):
            continue
        if preserve_decorations:
            has_text_frame = getattr(shape, "has_text_frame", False)
            if not has_text_frame:
                continue
        shape.element.getparent().remove(shape.element)


def is_title_like_shape(shape):
    name = getattr(shape, "name", "").lower()
    text = (shape.text or "").strip().lower() if getattr(shape, "has_text_frame", False) else ""
    top = emu_to_in(getattr(shape, "top", 9999999))
    return (
        "标题" in name
        or "title" in name
        or "标题" in text
        or "项目背景" in text
        or (top < 1.2 and bool(text) and len(text) <= 40)
        or text in ("xxx", "xxxx", "×××", "这是一个标题")
    )


def is_subtitle_like_shape(shape):
    name = getattr(shape, "name", "").lower()
    text = (shape.text or "").strip().lower() if getattr(shape, "has_text_frame", False) else ""
    return "副标题" in name or "subtitle" in name or "小标题" in text or text == "xxxxx"


def prepare_content_template(slide, slide_data):
    page_title = slide_data.get("page_title") or slide_data.get("title") or ""
    subtitle = slide_data.get("subtitle") or ""
    text_frames = all_text_frame_shapes(slide)
    title_shape = next((shape for shape in text_frames if is_title_like_shape(shape)), None)
    subtitle_shape = next((shape for shape in text_frames if shape is not title_shape and is_subtitle_like_shape(shape)), None)

    preserved = []
    if title_shape is not None:
        replace_text_preserve_shape(title_shape, page_title)
        preserved.append(title_shape.element)
    if subtitle_shape is not None:
        replace_text_preserve_shape(subtitle_shape, subtitle)
        preserved.append(subtitle_shape.element)

    for shape in list(all_text_frame_shapes(slide)):
        if any(shape.element is element for element in preserved):
            continue
        if not (shape.text or "").strip():
            continue
        shape.element.getparent().remove(shape.element)

    if title_shape is None:
        return None
    return emu_to_in(title_shape.top + title_shape.height)


def duplicate_slide(prs, source_slide):
    try:
        layout = source_slide.slide_layout
    except Exception:
        layout = prs.slide_layouts[0]
    dest = prs.slides.add_slide(layout)
    for shape in list(dest.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in source_slide.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dest.part.rels.get_or_add(rel.reltype, rel.target_part)
    return dest


def move_third_slide_last(prs):
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    if len(ids) >= 3:
        third = ids[2]
        sld_id_lst.remove(third)
        sld_id_lst.append(third)


def delete_slide(prs, index):
    sld_id_lst = prs.slides._sldIdLst
    slide_id = list(sld_id_lst)[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    sld_id_lst.remove(slide_id)


def move_slide(prs, old_index, new_index):
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    slide_id = slides[old_index]
    sld_id_lst.remove(slide_id)
    sld_id_lst.insert(new_index, slide_id)


def get_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            texts.append(shape.text.strip())
    return " ".join(texts).lower()


def looks_like_toc_text(text):
    if any(keyword in text for keyword in ["目录", "contents", "agenda"]):
        return True
    numbered_patterns = [
        r"\b0?[1-9][\.\、\)]",
        r"[一二三四五六七八九十][\.\、]",
    ]
    return any(len(re.findall(pattern, text)) >= 2 for pattern in numbered_patterns)


def detect_template_roles(prs):
    slide_count = len(prs.slides)
    if slide_count < 3:
        raise ValueError("style_clone expects at least 3 slides: cover, content style, closing.")

    cover_idx = 0
    end_idx = slide_count - 1
    toc_idx = None
    for idx in range(1, end_idx):
        text = get_slide_text(prs.slides[idx])
        if looks_like_toc_text(text):
            toc_idx = idx
            break

    content_indices = [
        idx for idx in range(1, end_idx)
        if idx != toc_idx
    ]
    if not content_indices:
        content_indices = [1]

    return cover_idx, toc_idx, content_indices, end_idx


def slide_file_index(slide_file):
    match = re.search(r"slide(\d+)\.xml", str(slide_file or ""))
    return int(match.group(1)) - 1 if match else None


def load_template_role_spec(template_path):
    spec_dir = template_path.parent
    for spec_path in spec_dir.glob("*_spec_corrected.json"):
        try:
            spec = json.loads(spec_path.read_text(encoding="utf-8-sig"))
        except Exception:
            continue
        if spec.get("template_file") != template_path.name:
            continue

        roles = {"cover": [], "toc": [], "transition": [], "content": [], "end": []}
        slide_types = spec.get("slide_types") or {}
        if not isinstance(slide_types, dict):
            return None

        for role, info in slide_types.items():
            if role == "content_styles" and isinstance(info, list):
                for item in info:
                    if not isinstance(item, dict):
                        continue
                    idx = slide_file_index(item.get("slide_file"))
                    if idx is not None:
                        roles["content"].append(idx)
                continue
            if not isinstance(info, dict):
                continue
            idx = slide_file_index(info.get("slide_file"))
            if idx is None:
                continue
            if role in ("content", "content_style"):
                roles["content"].append(idx)
            elif role in ("transition", "chapter"):
                roles["transition"].append(idx)
            elif role in roles:
                roles[role].append(idx)

        return roles
    return None


def resolve_template_roles(prs, template_path):
    spec_roles = load_template_role_spec(template_path)
    if spec_roles:
        slide_count = len(prs.slides)
        cover_idx = spec_roles["cover"][0] if spec_roles["cover"] else 0
        toc_idx = spec_roles["toc"][0] if spec_roles["toc"] else None
        content_indices = spec_roles["content"] or [
            idx for idx in range(1, slide_count - 1)
            if idx != toc_idx and idx not in spec_roles["transition"]
        ]
        end_idx = spec_roles["end"][0] if spec_roles["end"] else slide_count - 1
        transition_indices = spec_roles["transition"]
        return cover_idx, toc_idx, transition_indices, content_indices, end_idx

    cover_idx, toc_idx, content_indices, end_idx = detect_template_roles(prs)
    return cover_idx, toc_idx, [], content_indices, end_idx


def ensure_toc_slide(desired_slides, has_toc_template):
    if not has_toc_template:
        return [slide for slide in desired_slides if slide.get("type") != "toc"]
    if any(slide.get("type") == "toc" for slide in desired_slides):
        return desired_slides

    items = [
        slide.get("page_title") or slide.get("title")
        for slide in desired_slides
        if slide.get("type") == "content" and (slide.get("page_title") or slide.get("title"))
    ]
    if not items:
        return desired_slides

    slides = list(desired_slides)
    insert_at = 1 if slides and slides[0].get("type") == "cover" else 0
    slides.insert(insert_at, {"type": "toc", "title": "目录", "items": items[:8]})
    return slides


def replace_cover_text(slide, cover_data, content):
    shapes = text_shapes(slide)
    if not shapes:
        return

    title_text = cover_data.get("title") or content.get("title", "")
    date_text = cover_data.get("date") or content.get("date", "")
    subtitle_text = cover_data.get("subtitle") or content.get("subtitle", "")

    title_shape = max(shapes, key=lambda s: (s.width or 0) * (s.height or 0))
    replace_text_preserve_shape(title_shape, title_text)

    remaining = [s for s in shapes if s is not title_shape]
    date_shape = None
    if remaining and date_text:
        date_shape = max(remaining, key=lambda s: s.top or 0)
        replace_text_preserve_shape(date_shape, date_text)

    for shape in remaining:
        if shape is date_shape:
            continue
        original = shape.text.strip().lower()
        is_placeholder_text = any(token in original for token in ["xxx", "标题", "副标题", "subtitle", "tagline"])
        is_hash_like = len(original) > 80 and all(ch in "0123456789abcdef" for ch in original)
        if subtitle_text:
            replace_text_preserve_shape(shape, subtitle_text)
        elif is_placeholder_text or is_hash_like:
            replace_text_preserve_shape(shape, "")


def replace_end_text(slide, end_data, content):
    shapes = text_shapes(slide)
    if not shapes:
        return

    title_text = end_data.get("title") or content.get("end_title") or "谢谢"
    date_text = end_data.get("date") or content.get("date", "")

    title_shape = max(shapes, key=lambda s: (s.width or 0) * (s.height or 0))
    replace_text_preserve_shape(title_shape, title_text)

    remaining = [s for s in shapes if s is not title_shape]
    date_shape = None
    if remaining and date_text:
        date_shape = max(remaining, key=lambda s: s.top or 0)
        replace_text_preserve_shape(date_shape, date_text)

    for shape in remaining:
        if shape is date_shape:
            continue
        original = shape.text.strip().lower()
        is_placeholder_text = any(token in original for token in ["xxx", "标题", "副标题", "thanks", "thank you", "谢谢"])
        if is_placeholder_text:
            replace_text_preserve_shape(shape, "")


def stat(slide, x, y, w, num, label, note, color):
    rect(slide, x, y, w, 1.32)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(1.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color)
    bar.line.fill.background()
    tx(slide, num, x + 0.2, y + 0.16, w - 0.35, 0.38, 22, color, True)
    tx(slide, label, x + 0.22, y + 0.69, w - 0.35, 0.22, 10.8, TEXT, True)
    tx(slide, note, x + 0.22, y + 0.98, w - 0.35, 0.22, 8.1, MUTED)


def bullet(slide, x, y, w, head, body, color=BLUE):
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.13), Inches(0.13))
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(color)
    dot.line.fill.background()
    tx(slide, head, x + 0.22, y, w - 0.22, 0.19, 10.2, TEXT, True)
    tx(slide, body, x + 0.22, y + 0.25, w - 0.22, 0.42, 8.5, MUTED)


def layout_kpi(slide, items, start_y=1.35):
    colors = [BLUE, GREEN, CYAN, ORANGE]
    cols = min(len(items), 3)
    card_w = 3.55 if cols == 3 else 5.2
    start_x = 0.72 if cols == 3 else 1.15
    gap = 0.55
    for i, item in enumerate(items[:cols]):
        stat(
            slide,
            start_x + i * (card_w + gap),
            start_y,
            card_w,
            item.get("value", ""),
            item.get("label", ""),
            item.get("desc", ""),
            colors[i % len(colors)],
        )
    for i, item in enumerate(items[cols:cols + 4]):
        x = 0.82 + (i % 2) * 5.93
        y = start_y + 1.97 + (i // 2) * 1.05
        bullet(slide, x, y, 4.9, item.get("label", ""), item.get("desc", ""), colors[(i + cols) % len(colors)])


def layout_grid(slide, items, start_y=1.45):
    colors = [BLUE, CYAN, GREEN, ORANGE, BLUE, GREEN]
    for i, item in enumerate(items[:6]):
        x = 0.75 + (i % 2) * 6.05
        y = start_y + (i // 2) * 1.45
        rect(slide, x, y, 5.5, 1.02)
        tx(slide, f"{i + 1:02d}", x + 0.28, y + 0.28, 0.45, 0.15, 10.5, colors[i], True)
        tx(slide, item.get("label", ""), x + 0.9, y + 0.18, 2.6, 0.18, 10.5, TEXT, True)
        if item.get("value"):
            tx(slide, item["value"], x + 3.8, y + 0.18, 1.2, 0.18, 10.5, colors[i], True, "right")
        tx(slide, item.get("desc", ""), x + 0.9, y + 0.52, 3.95, 0.32, 8.4, MUTED)


def layout_timeline(slide, items, start_y=1.45):
    colors = [BLUE, GREEN, CYAN, ORANGE, BLUE]
    line_y = start_y + 2.4
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.4), Inches(line_y), Inches(10.2), Inches(0.03)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(LIGHT_BLUE)
    slide.shapes[-1].line.fill.background()
    step_w = 10.4 / max(1, len(items))
    for i, item in enumerate(items[:5]):
        cx = 1.5 + i * step_w + step_w / 2
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.22), Inches(line_y - 0.2), Inches(0.44), Inches(0.44))
        node.fill.solid()
        node.fill.fore_color.rgb = rgb(colors[i])
        node.line.fill.background()
        tx(slide, str(i + 1), cx - 0.2, line_y - 0.12, 0.4, 0.12, 8.5, "FFFFFF", True, "center")
        tx(slide, item.get("label", ""), cx - step_w / 2 + 0.05, start_y + 0.93, step_w - 0.1, 0.2, 10.2, TEXT, True, "center")
        tx(slide, item.get("desc", ""), cx - step_w / 2 + 0.1, 4.28, step_w - 0.2, 0.6, 8.2, MUTED, False, "center")


def layout_big(slide, items, start_y=1.75):
    item = items[0] if items else {}
    rect(slide, 1.05, start_y, 11.0, 3.7)
    tx(slide, item.get("value", ""), 1.45, start_y + 0.3, 10.2, 0.75, 40, BLUE, True, "center")
    tx(slide, item.get("label", ""), 1.45, start_y + 1.3, 10.2, 0.35, 18, TEXT, True, "center")
    tx(slide, item.get("desc", ""), 1.75, start_y + 2.0, 9.6, 0.6, 12, MUTED, False, "center")


def render_toc_slide(slide, slide_data, page_num, source):
    items = slide_data.get("items", [])
    if not items:
        return

    def is_toc_marker(text):
        normalized = text.strip().lower()
        if normalized in ("目录", "contents", "content", "agenda", "part"):
            return True
        return bool(re.fullmatch(r"(?:0?[1-9]|part\s*0?[1-9])", normalized))

    candidates = []
    for shape in iter_text_shapes(slide.shapes):
        raw = shape.text.strip()
        normalized = raw.lower()
        if is_toc_marker(raw):
            continue
        if (
            "×" in raw
            or "x" in normalized
            or "标题" in raw
            or "部分" in raw
            or re.match(r"^[一二三四五六七八九十][、.]", raw)
            or re.match(r"^0?[1-9][、. )]", raw)
        ):
            candidates.append(shape)

    if len(candidates) < len(items):
        extra = [
            shape
            for shape in iter_text_shapes(slide.shapes)
            if shape not in candidates
            and not is_toc_marker(shape.text)
        ]
        candidates.extend(extra)

    def replacement(original, item):
        text = original.strip()
        match = re.match(r"^([一二三四五六七八九十][、.]\s*)(.*)$", text)
        if match:
            return match.group(1) + str(item)
        match = re.match(r"^(0?[1-9][、. )]\s*)(.*)$", text)
        if match:
            return match.group(1) + str(item)
        return str(item)

    if len(candidates) == 1 and len(items) > 1:
        original_lines = [line.strip() for line in candidates[0].text.splitlines() if line.strip()]
        prefixes = []
        for line in original_lines:
            match = re.match(r"^([一二三四五六七八九十][、.]\s*)", line)
            if not match:
                match = re.match(r"^(0?[1-9][、. )]\s*)", line)
            prefixes.append(match.group(1) if match else "")
        if not prefixes:
            prefixes = [""] * len(items)
        while len(prefixes) < len(items):
            prefixes.append(f"{len(prefixes) + 1:02d} ")
        joined = "\n\n".join(prefixes[idx] + str(item) for idx, item in enumerate(items))
        replace_text_preserve_shape(candidates[0], joined)
        return

    for idx, item in enumerate(items):
        if idx >= len(candidates):
            break
        shape = candidates[idx]
        replace_text_preserve_shape(shape, replacement(shape.text, item))

    for idx in range(len(items), len(candidates)):
        replace_text_preserve_shape(candidates[idx], "")


def render_transition_slide(slide, slide_data, page_num, source):
    chapter = slide_data.get("chapter") or slide_data.get("section") or ""
    heading = slide_data.get("title") or slide_data.get("page_title") or ""
    text_shapes_in_order = list(iter_text_shapes(slide.shapes))
    if not text_shapes_in_order:
        return

    if len(text_shapes_in_order) == 1:
        original = text_shapes_in_order[0].text.strip()
        if chapter:
            match = re.match(r"^(\S+)(\s+)(.*)$", original)
            if match:
                replace_text_preserve_shape(text_shapes_in_order[0], f"{chapter}{match.group(2)}{heading}")
            else:
                replace_text_preserve_shape(text_shapes_in_order[0], f"{chapter}   {heading}")
        else:
            replace_text_preserve_shape(text_shapes_in_order[0], heading)
        return

    if chapter:
        replace_text_preserve_shape(text_shapes_in_order[0], chapter)
        replace_text_preserve_shape(text_shapes_in_order[1], heading)
    else:
        replace_text_preserve_shape(text_shapes_in_order[0], heading)


def render_content_slide(slide, slide_data, page_num, source, draw_title=True, content_top=None):
    if draw_title:
        title(slide, slide_data.get("page_title") or slide_data.get("title", ""), slide_data.get("subtitle"))
    start_y = max(1.35, (content_top or 1.1) + 0.28)
    body = slide_data.get("body", {})
    items = body.get("items", [])
    layout = body.get("layout", "auto")
    if layout == "auto":
        layout = "kpi" if any(item.get("value") for item in items) else ("timeline" if len(items) == 4 and body.get("timeline") else "grid")
    if layout in ("kpi", "kpi_card"):
        layout_kpi(slide, items, start_y)
    elif layout == "timeline":
        layout_timeline(slide, items, start_y)
    elif layout in ("big_stat", "quote"):
        layout_big(slide, items, max(1.75, start_y))
    else:
        layout_grid(slide, items, start_y)
    footer(slide, page_num, source)


def build(content_path: Path, output_path: Path, template_path: Path):
    content = json.loads(content_path.read_text(encoding="utf-8-sig"))
    prs = Presentation(str(template_path))

    cover_idx, toc_template_idx, transition_indices, content_indices, end_idx = resolve_template_roles(prs, template_path)
    desired_slides = ensure_toc_slide(content.get("slides") or [], toc_template_idx is not None)
    body_slides = [s for s in desired_slides if s.get("type") in ("toc", "transition", "chapter", "content")]
    toc_source = prs.slides[toc_template_idx] if toc_template_idx is not None else prs.slides[content_indices[0]]
    transition_sources = [prs.slides[idx] for idx in transition_indices] or [prs.slides[content_indices[0]]]
    content_sources = [prs.slides[idx] for idx in content_indices]
    end_source = prs.slides[end_idx]
    style_selection = content.get("style_selection", "rotate")
    rng = random.Random(content.get("style_seed", 20260430))

    generated = []
    content_style_cursor = 0
    transition_style_cursor = 0
    for slide_data in body_slides:
        if slide_data.get("type") == "toc":
            generated.append((duplicate_slide(prs, toc_source), slide_data))
        elif slide_data.get("type") in ("transition", "chapter"):
            source = transition_sources[transition_style_cursor % len(transition_sources)]
            transition_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))
        else:
            if style_selection == "random":
                source = rng.choice(content_sources)
            else:
                source = content_sources[content_style_cursor % len(content_sources)]
                content_style_cursor += 1
            generated.append((duplicate_slide(prs, source), slide_data))

    generated_end = end_source

    # Keep only original cover plus generated body/end slides.
    keep_ids = {prs.slides[cover_idx].slide_id}
    keep_ids.update(slide.slide_id for slide, _ in generated)
    keep_ids.add(generated_end.slide_id)
    for idx in range(len(prs.slides) - 1, -1, -1):
        if prs.slides[idx].slide_id not in keep_ids:
            delete_slide(prs, idx)

    # Reorder to cover, generated body slides, generated end.
    target_order = [prs.slides[0].slide_id] + [slide.slide_id for slide, _ in generated] + [generated_end.slide_id]
    for target_idx, slide_id in enumerate(target_order):
        current_idx = next(i for i, slide in enumerate(prs.slides) if slide.slide_id == slide_id)
        if current_idx != target_idx:
            move_slide(prs, current_idx, target_idx)

    slides = list(prs.slides)
    template_title_bottom = {}
    for slide, slide_data in generated:
        if slide_data.get("type") in ("toc", "transition", "chapter"):
            continue
        template_title_bottom[slide.slide_id] = prepare_content_template(slide, slide_data)

    cover_data = next((s for s in desired_slides if s.get("type") == "cover"), {})
    replace_cover_text(slides[0], cover_data, content)

    source = content.get("source", "")
    for i, slide_data in enumerate(body_slides, start=1):
        if slide_data.get("type") == "toc":
            render_toc_slide(slides[i], slide_data, i + 1, source)
        elif slide_data.get("type") in ("transition", "chapter"):
            render_transition_slide(slides[i], slide_data, i + 1, source)
        else:
            render_content_slide(
                slides[i],
                slide_data,
                i + 1,
                source,
                draw_title=template_title_bottom.get(slides[i].slide_id) is None,
                content_top=template_title_bottom.get(slides[i].slide_id),
            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(output_path)


def main():
    if len(sys.argv) < 3:
        print("Usage: python render_geely_report.py content.json output.pptx [template.pptx]")
        sys.exit(1)
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    template = Path(sys.argv[3]) if len(sys.argv) > 3 else skill_dir / "templates" / "吉利模版.pptx"
    build(Path(sys.argv[1]), Path(sys.argv[2]), template)


if __name__ == "__main__":
    main()
